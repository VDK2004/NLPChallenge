import uuid
from typing import List, Dict, Union
from llama_index import (
    SimpleDirectoryReader,
    Document,
    ServiceContext,
    VectorStoreIndex
)
from llama_index.node_parser import SimpleNodeParser
from llama_index.extractors import TitleExtractor
from llama_index.readers.youtube_transcript import YoutubeTranscriptReader
from pytube import YouTube
import os
import tempfile
from bs4 import BeautifulSoup
import requests
from pptx import Presentation

class DocumentProcessor:
    def __init__(self):
        self.processed_docs = []
        self.node_parser = SimpleNodeParser.from_defaults(
            chunk_size=1024,
            chunk_overlap=20
        )
        self.title_extractor = TitleExtractor(nodes=5)

    def process_pdf(self, file) -> Document:
        """Process a PDF file using LlamaIndex"""
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
            temp_file.write(file.getvalue())
            temp_path = temp_file.name

        try:
            # Use LlamaIndex's SimpleDirectoryReader for PDF
            reader = SimpleDirectoryReader(input_files=[temp_path])
            docs = reader.load_data()
            
            # Add metadata
            for doc in docs:
                doc.metadata.update({
                    'id': str(uuid.uuid4()),
                    'filename': file.name,
                    'source_type': 'pdf'
                })
            
            return docs
        finally:
            os.unlink(temp_path)

    def process_youtube(self, url: str) -> List[Document]:
        """Process a YouTube video"""
        loader = YoutubeTranscriptReader()
        
        # Extract video ID from URL
        video_id = url.split('v=')[-1].split('&')[0]
        
        try:
            # Try to get video info
            yt = YouTube(url)
            title = yt.title
        except Exception:
            # Fallback to using video ID as title if we can't get the actual title
            title = f"YouTube Video ({video_id})"
        
        # Load transcript
        docs = loader.load_data(ytlinks=[url])
        
        # Add metadata
        for doc in docs:
            doc.metadata.update({
                'id': str(uuid.uuid4()),
                'filename': title,
                'source_type': 'youtube',
                'url': url,
                'video_id': video_id
            })
        
        return docs

    def process_webpage(self, url: str) -> List[Document]:
        """Process a webpage"""
        response = requests.get(url)
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Extract main content
        content = soup.get_text()
        
        # Create document
        doc = Document(
            text=content,
            metadata={
                'id': str(uuid.uuid4()),
                'filename': url,
                'source_type': 'webpage',
                'url': url
            }
        )
        
        return [doc]

    def process_powerpoint(self, file) -> List[Document]:
        """Process a PowerPoint file"""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_file.write(file.getvalue())
            temp_path = temp_file.name

        try:
            prs = Presentation(temp_path)
            content = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                content.append("\n".join(slide_text))
            
            doc = Document(
                text="\n\n".join(content),
                metadata={
                    'id': str(uuid.uuid4()),
                    'filename': file.name,
                    'source_type': 'powerpoint'
                }
            )
            
            return [doc]
        finally:
            os.unlink(temp_path)

    def process_documents(self, files: List[Union[str, tempfile._TemporaryFileWrapper]]) -> List[Dict]:
        """
        Process multiple documents and return their content with metadata
        Supports PDF files, YouTube URLs, webpages, and PowerPoint files
        """
        all_docs = []
        
        for file in files:
            try:
                if isinstance(file, str) and ('youtube.com' in file or 'youtu.be' in file):
                    docs = self.process_youtube(file)
                elif isinstance(file, str) and (file.startswith('http://') or file.startswith('https://')):
                    docs = self.process_webpage(file)
                else:
                    if file.name.lower().endswith('.pdf'):
                        docs = self.process_pdf(file)
                    elif file.name.lower().endswith(('.pptx', '.ppt')):
                        docs = self.process_powerpoint(file)
                    else:
                        continue
                
                # Convert LlamaIndex documents to our format
                for doc in docs:
                    doc_info = {
                        'id': doc.metadata['id'],
                        'filename': doc.metadata['filename'],
                        'content': doc.text,
                        'source_type': doc.metadata.get('source_type', 'unknown'),
                        'num_pages': len(self.node_parser.get_nodes_from_documents([doc]))
                    }
                    
                    # Add URL for web sources
                    if 'url' in doc.metadata:
                        doc_info['url'] = doc.metadata['url']
                    
                    all_docs.append(doc_info)
            
            except Exception as e:
                raise Exception(f"Error processing document {getattr(file, 'name', file)}: {str(e)}")
        
        return all_docs
